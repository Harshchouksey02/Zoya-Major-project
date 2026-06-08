from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

# Color Scheme
DARK_BLUE = RGBColor(27, 54, 93)     # Primary Header / Backgrounds
LIGHT_BLUE = RGBColor(74, 144, 226)   # Accent Highlights
DARK_GRAY = RGBColor(48, 48, 48)      # Body Text
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(245, 246, 248)
EMERALD_GREEN = RGBColor(22, 163, 74)

blank_layout = prs.slide_layouts[6]

# Helper: Add title slide background and content
def build_title_slide():
    slide = prs.slides.add_slide(blank_layout)
    
    # Fill background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BLUE
    bg.line.color.rgb = DARK_BLUE

    # Top border line
    top_border = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.15))
    top_border.fill.solid()
    top_border.fill.fore_color.rgb = LIGHT_BLUE
    top_border.line.color.rgb = LIGHT_BLUE

    # Title box
    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(1.2), Inches(11.33), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = "AgroVeda"
    p.font.name = 'Georgia'
    p.font.size = Pt(64)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    p2 = tf.add_paragraph()
    p2.text = "E-Commerce Platform for Agricultural Products"
    p2.font.name = 'Arial'
    p2.font.size = Pt(28)
    p2.font.color.rgb = LIGHT_BLUE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    # Submission box (Left)
    sub_box = slide.shapes.add_textbox(Inches(1.5), Inches(3.8), Inches(4.5), Inches(2.2))
    tf_sub = sub_box.text_frame
    tf_sub.word_wrap = True
    p_sub = tf_sub.paragraphs[0]
    p_sub.text = "Submitted By:"
    p_sub.font.name = 'Arial'
    p_sub.font.size = Pt(16)
    p_sub.font.bold = True
    p_sub.font.color.rgb = LIGHT_BLUE
    
    student_details = [
        "ZOYA IBRAHIM",
        "Enrollment No: 0210CA241134",
        "4th Semester MCA | Session 2025-26"
    ]
    for d in student_details:
        p_d = tf_sub.add_paragraph()
        p_d.text = d
        p_d.font.name = 'Arial'
        p_d.font.size = Pt(14)
        p_d.font.color.rgb = WHITE
        p_d.space_before = Pt(5)

    # Guidance box (Right)
    guide_box = slide.shapes.add_textbox(Inches(7.33), Inches(3.8), Inches(4.5), Inches(2.2))
    tf_guide = guide_box.text_frame
    tf_guide.word_wrap = True
    p_guide = tf_guide.paragraphs[0]
    p_guide.text = "Under the Guidance of:"
    p_guide.font.name = 'Arial'
    p_guide.font.size = Pt(16)
    p_guide.font.bold = True
    p_guide.font.color.rgb = LIGHT_BLUE
    
    guide_details = [
        "Prof. Sachin Dubey",
        "Department of MCA",
        "Shri Ram Institute of Technology, Jabalpur"
    ]
    for gd in guide_details:
        p_gd = tf_guide.add_paragraph()
        p_gd.text = gd
        p_gd.font.name = 'Arial'
        p_gd.font.size = Pt(14)
        p_gd.font.color.rgb = WHITE
        p_gd.space_before = Pt(5)

# Helper: Add Standard Header
def add_slide_header(slide, title_text):
    # Banner
    banner = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.1))
    banner.fill.solid()
    banner.fill.fore_color.rgb = DARK_BLUE
    banner.line.color.rgb = DARK_BLUE
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(0.18), Inches(12.13), Inches(0.75))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = 'Arial'
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE
    
    # Accent Line
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.33), Inches(0.08))
    accent.fill.solid()
    accent.fill.fore_color.rgb = LIGHT_BLUE
    accent.line.color.rgb = LIGHT_BLUE

    # Footer
    add_slide_footer(slide)

def add_slide_footer(slide):
    footer_text = "AgroVeda MCA Major Project"
    txBox = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.13), Inches(0.4))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = footer_text
    p.font.name = 'Arial'
    p.font.size = Pt(10)
    p.font.color.rgb = RGBColor(160, 160, 160)
    p.alignment = PP_ALIGN.RIGHT

# Helper: Add Bullet points
def add_bullet_points(slide, points, left=0.8, top=1.8, width=11.7, height=4.8, font_size=18):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + pt
        p.font.name = 'Arial'
        p.font.size = Pt(font_size)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(14)

# Helper: Add double columns
def add_two_column_text(slide, col1_title, col1_points, col2_title, col2_points):
    # Col 1
    box1 = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
    tf1 = box1.text_frame
    tf1.word_wrap = True
    p1 = tf1.paragraphs[0]
    p1.text = col1_title
    p1.font.name = 'Arial'
    p1.font.size = Pt(20)
    p1.font.bold = True
    p1.font.color.rgb = DARK_BLUE
    p1.space_after = Pt(12)
    for pt in col1_points:
        p = tf1.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

    # Col 2
    box2 = slide.shapes.add_textbox(Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
    tf2 = box2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = col2_title
    p2.font.name = 'Arial'
    p2.font.size = Pt(20)
    p2.font.bold = True
    p2.font.color.rgb = DARK_BLUE
    p2.space_after = Pt(12)
    for pt in col2_points:
        p = tf2.add_paragraph()
        p.text = "• " + pt
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(10)

# Helper: Add image with description
def add_image_slide(slide, image_path, description_points, image_width=6.0, img_height=4.5):
    # Add Image
    try:
        if os.path.exists(image_path):
            slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.8), width=Inches(image_width), height=Inches(img_height))
            left_offset = 1.2 + image_width
        else:
            # Placeholder box
            placeholder = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(image_width), Inches(img_height))
            placeholder.fill.solid()
            placeholder.fill.fore_color.rgb = LIGHT_GRAY
            placeholder.line.color.rgb = LIGHT_BLUE
            tx = placeholder.text_frame
            p = tx.paragraphs[0]
            p.text = f"[Image Placeholder: {os.path.basename(image_path)}]"
            p.alignment = PP_ALIGN.CENTER
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            left_offset = 1.2 + image_width
    except Exception as e:
        print(f"Error adding image {image_path}: {e}")
        left_offset = 7.0

    # Add Points
    txBox = slide.shapes.add_textbox(Inches(left_offset), Inches(1.8), Inches(12.5 - left_offset), Inches(img_height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, pt in enumerate(description_points):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = "• " + pt
        p.font.name = 'Arial'
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)

# ============================================================
# BUILD SLIDES
# ============================================================

# 1. Title Slide
build_title_slide()

# 2. Agenda / Outline
slide2 = prs.slides.add_slide(blank_layout)
add_slide_header(slide2, "Agenda")
agenda_points = [
    "Project Introduction & Domain Overview",
    "Problem Statement & Project Objectives",
    "Analysis of Existing vs. Proposed System",
    "System Architecture, Workflows & Tech Stack",
    "Database Design & UML Activity Modeling",
    "Core Code Algorithms & Security Implementations",
    "Admin Order & Sales Registry (Dashboard Demonstration)",
    "Advantages, Future Scope & Conclusions"
]
add_bullet_points(slide2, agenda_points, font_size=20)

# 3. Project Introduction
slide3 = prs.slides.add_slide(blank_layout)
add_slide_header(slide3, "Project Introduction")
intro_points = [
    "Agriculture is the backbone of the Indian economy, employing over 50% of the workforce, but digital input procurement remains highly unstructured.",
    "Traditional channels rely on local physical retail dealers, leading to long travel distances, lack of inventory selection, and high pricing markups.",
    "AgroVeda is a full-stack agricultural e-commerce platform designed to bridge this digital divide.",
    "Enables farmers and cooperatives to purchase quality pesticides, fertilizers, growth regulators, and crop boosters directly from their mobile devices.",
    "Designed with a dual-language (Hindi/English) catalog interface to ensure high accessibility for rural farmers."
]
add_bullet_points(slide3, intro_points, font_size=16)

# 4. Problem Statement
slide4 = prs.slides.add_slide(blank_layout)
add_slide_header(slide4, "Problem Statement")
prob_points = [
    "Middlemen Inefficiencies: Complex retail dealership networks result in high markup prices for vital chemicals/boosters.",
    "Asymmetric Information: Farmers lack transparent pricing and have no direct way to compare product ingredients and specifications.",
    "Language Barrier: Most commercial e-commerce portals are in English, excluding the major demographic of Hindi-speaking rural farmers.",
    "Logistical Overhead: Farmers travel multiple kilometers to towns without guarantee of stock availability, leading to wasted time and labor.",
    "Trust & Payment Issues: Absence of cash-on-delivery options or simplified, secure local digital billing gateways (e.g. Razorpay UPI QR simulation)."
]
add_bullet_points(slide4, prob_points, font_size=16)

# 5. Objectives
slide5 = prs.slides.add_slide(blank_layout)
add_slide_header(slide5, "Project Objectives")
obj_points = [
    "Develop a responsive, full-stack, mobile-friendly agricultural input e-commerce portal.",
    "Implement a bilingual interface (English and Hindi) for seamless navigation and product catalogs.",
    "Create a secure, JWT-authenticated User Registration and Login system.",
    "Integrate Razorpay payment gateway simulation for UPI, debit/credit cards, alongside a Cash on Delivery (COD) workflow.",
    "Construct a secure Admin Dashboard enabling product CRUD operations (Add, delete, update pricing) and a real-time Orders & Sales Summary monitor to track who purchased, order timestamps, and items summary."
]
add_bullet_points(slide5, obj_points, font_size=16)

# 6. Existing System vs Drawbacks
slide6 = prs.slides.add_slide(blank_layout)
add_slide_header(slide6, "Existing System & Drawbacks")
add_two_column_text(
    slide6,
    "Current Traditional System",
    [
        "Procurement through local village or town shops.",
        "Rely on word-of-mouth recommendations from shop owners.",
        "Manual bookkeeping (paper invoices) with no transaction histories.",
        "Cash-only transactions with high markup commissions."
    ],
    "Key Critical Drawbacks",
    [
        "Limited catalog variety based on shop keeper's preferences.",
        "Lack of pricing transparency (varying prices across different shops).",
        "Wasted labor and transportation costs traveling to local dealers.",
        "No digital bills, order tracking, or verification receipts."
    ]
)

# 7. Proposed System
slide7 = prs.slides.add_slide(blank_layout)
add_slide_header(slide7, "Proposed System Solution")
prop_points = [
    "Centralized E-Commerce Hub: Direct platform connecting farmers with catalog listings.",
    "Hindi/English Support: Easy-toggle language context supporting titles and descriptions.",
    "Secure Cart & Order Lifecycles: Full workflow covering add-to-cart, quantity updates, order summary validations, and automated billing details.",
    "Simulated Payment Gateways: Integration with Razorpay API (with instant QR scan/simulate-success testing modal).",
    "Interactive Admin CRM Dashboard: Real-time Order Registry logging customer mobile, timestamp of purchase, shipping address, total sales revenue statistics, and order items summary."
]
add_bullet_points(slide7, prop_points, font_size=16)

# 8. System Architecture
slide8 = prs.slides.add_slide(blank_layout)
add_slide_header(slide8, "System Architecture")
arch_points = [
    "Model-View-Controller (MVC) architectural design pattern deployed.",
    "Client-Side View: Single Page Application (SPA) built using React 18, TypeScript, and Vite for fast builds and components render.",
    "Server-Side Controller: Express.js (Node.js) serving REST endpoints (/api/products, /api/orders, /api/cart_items) deployed as serverless functions.",
    "Database Layer: Relational PostgreSQL database managed with an active connection pool, utilizing a local file-based JSON database engine fallback for offline local testing.",
    "Security Shield: Cross-Origin Resource Sharing (CORS) configurations, bcrypt password salting, and JSON Web Token (JWT) header authorization check."
]
add_image_slide(slide8, "diagrams/dfd_level0.png", arch_points, image_width=6.0)

# 9. Methodology / Workflow
slide9 = prs.slides.add_slide(blank_layout)
add_slide_header(slide9, "Methodology & Development Workflow")
meth_points = [
    "Agile Scrum Methodology implemented, divided into 6 distinct sprints over 20 days:",
    "Sprint 1 (Days 1-3): Project Setup & DB Schema Design. Setup Vite frontend scaffolding and local Express API route definitions.",
    "Sprint 2 (Days 4-7): UI Design & Bilingual Integration. Implemented homepage slides, product list filters, and English/Hindi language toggle.",
    "Sprint 3 (Days 8-10): Auth Shield. Built JWT signup, login, password encryption, and role-based route guards.",
    "Sprint 4 (Days 11-14): Cart & Checkout. Added shopping cart store, quantity managers, and address validation forms.",
    "Sprint 5 (Days 15-17): Razorpay Payments. Deployed backend order generation and signature HMAC-SHA256 checking.",
    "Sprint 6 (Days 18-20): Admin Panel, Orders Monitor & Vercel deployment."
]
add_bullet_points(slide9, meth_points, font_size=14)

# 10. Technologies Used
slide10 = prs.slides.add_slide(blank_layout)
add_slide_header(slide10, "Technologies Used")
add_two_column_text(
    slide10,
    "Frontend Stack",
    [
        "React 18 & TypeScript (Type-safe components)",
        "Vite (Next-generation bundle tool)",
        "TailwindCSS & Shadcn UI (Custom layout styling)",
        "React Router DOM (Routing)",
        "Lucide React (Modern iconography)"
    ],
    "Backend & Database Stack",
    [
        "Node.js & Express.js (REST API server)",
        "PostgreSQL (ACID-compliant Database)",
        "JSON Web Tokens (JWT) (Authentication)",
        "Bcrypt.js (Password encryption)",
        "Razorpay SDK (Payments gateway integration)",
        "Vercel (Production hosting platform)"
    ]
)

# 11. UML Activity Diagram
slide11 = prs.slides.add_slide(blank_layout)
add_slide_header(slide11, "UML Diagrams - Activity Diagram")
uml_points = [
    "This Activity Diagram outlines the execution flow of the Customer checkout and order confirmation process.",
    "Flow triggers on 'Checkout' request.",
    "Validates stock availability by querying Products DB.",
    "Fork Decision: Cash on Delivery (COD) vs. Razorpay Online Payment.",
    "Online payment initiates Razorpay API Order token creation and launches the UI payment widget.",
    "Upon signature verification, order data is committed to the database, cart items are cleared, and invoice receipt notifications are fired."
]
add_image_slide(slide11, "diagrams/activity_diagram.png", uml_points, image_width=6.0)

# 12. Database Design & Schema
slide12 = prs.slides.add_slide(blank_layout)
add_slide_header(slide12, "Database Design (ER Diagram)")
db_points = [
    "Our database schema consists of six core relational tables:",
    "USERS: Holds registration emails, passwords hashes, and metadata.",
    "USER_ROLES: Handles role permissions (admin vs. customer).",
    "PRODUCTS: Holds product detail metadata (prices, units, usage, names in Hindi/English).",
    "CART_ITEMS: Maps user shopping baskets (product reference, quantities).",
    "ORDERS: Logs final checkout transactions (customer phone, address, timestamp, status, amounts).",
    "ORDER_ITEMS: Captures product snapshot details inside orders."
]
add_image_slide(slide12, "diagrams/er_diagram.png", db_points, image_width=6.0)

# 13. Algorithm / Model (Security & Verify)
slide13 = prs.slides.add_slide(blank_layout)
add_slide_header(slide13, "Core Backend Algorithms")
alg_points = [
    "1. Authentication & Security (JWT Flow):",
    "   - Customer password hashed on signup using bcryptjs with a work factor of 10.",
    "   - On successful login, server creates a JWT payload containing { userId, email, role } signed using HMAC-SHA256 with a private key, returning an authorization token.",
    "   - Auth Middleware intercepts private routes, decodes headers Bearer tokens, and sets context parameters.",
    "2. Razorpay Signature Verification Algorithm:",
    "   - When Razorpay completes a transaction, it returns: razorpay_order_id, razorpay_payment_id, and razorpay_signature.",
    "   - The server computes the expected signature using HMAC-SHA256:",
    "     signature = HMAC_SHA256(order_id + '|' + payment_id, secret_key)",
    "   - If expected signature matches razorpay_signature, order is verified."
]
add_bullet_points(slide13, alg_points, font_size=14)

# 14. Screenshots / Demonstration
slide14 = prs.slides.add_slide(blank_layout)
add_slide_header(slide14, "User Interface Modules")
demo_points = [
    "Homepage Module: Features a bilingual sliding banner showing seasonal recommendations, and a category grid to filter fertilizers and fungicides.",
    "Product Catalogue: Dual-language support toggle. Displays crop usage descriptions, ratings, pricing, and bulk offers (e.g. Buy 5 Get 1 Free).",
    "Cart & Checkout: Address entry validation using Zod schemas, mobile validation (Indian mobile structure check), and payment mode selection.",
    "Order Confirmation Screen: Simulates Razorpay transaction receipt rendering invoice items and details."
]
add_image_slide(slide14, "diagrams/use_case_diagram.png", demo_points, image_width=6.0)

# 15. Admin Panel (Order Summary & Sales Registry)
slide15 = prs.slides.add_slide(blank_layout)
add_slide_header(slide15, "Admin Order & Sales Registry")
admin_points = [
    "Admin Dashboard now supports a dual-panel layout:",
    "1. Product Inventory Manager: Admin can create new items (names, category, price, Hindi description) and delete existing ones.",
    "2. Order & Sales Summary Registry: Displays key sales statistics (Total Revenue, total orders counts, completed/COD counts).",
    "3. Full Buy Summary: Lists all orders dynamically. Each record contains:",
    "   - Order number and purchase timestamp (formatted to Indian Standard Time).",
    "   - Customer credentials: Name, Phone Number, and complete shipping address.",
    "   - Detailed bill calculations: Subtotals, discount deductions, and final paid values.",
    "   - Interactive Order Status manager dropdown (Change order from pending -> processing -> delivered -> cancelled)."
]
add_bullet_points(slide15, admin_points, font_size=15)

# 16. Advantages & Future Scope
slide16 = prs.slides.add_slide(blank_layout)
add_slide_header(slide16, "Advantages & Future Scope")
add_two_column_text(
    slide16,
    "Key Advantages",
    [
        "Direct Procurement: Removes retail agent layers, reducing input purchase cost by 15-20%.",
        "Accessibility: Simple bilingual layout works smoothly on mobile browsers.",
        "Transparency: Fixed pricing with instant invoice summaries.",
        "Security: Safe role-based admin panels and verified payment gateway check."
    ],
    "Future Scope & Enhancements",
    [
        "AI Crop Advisor: Recommendation engine based on local soil health test details.",
        "Voice Search: Multi-lingual voice recognition queries for illiterate farmers.",
        "WhatsApp API Integration: Automated dispatch alerts sent directly to phone numbers.",
        "Multi-Vendor Marketplace: Allow verified local manufacturers to sell directly."
    ]
)

# 17. Conclusion & References
slide17 = prs.slides.add_slide(blank_layout)
add_slide_header(slide17, "Conclusion & References")
ref_points = [
    "Conclusion:",
    "   - AgroVeda successfully delivers a scalable, type-safe e-commerce solution for agricultural distribution.",
    "   - Demonstrated the practical application of React, Node.js, Express, and PostgreSQL/In-Memory DB.",
    "   - Bridged the rural digital gap using bilingual catalog design and interactive Razorpay/COD flows.",
    "   - Empowered administrators with dynamic sales tracking and order status monitoring dashboards.",
    "References:",
    "   - React 18 and Vite Documentation: https://react.dev | https://vite.dev",
    "   - Express.js Node Routing API Reference: https://expressjs.com",
    "   - PostgreSQL client driver documentation: https://node-postgres.com",
    "   - Razorpay Web checkout integration APIs: https://razorpay.com/docs"
]
add_bullet_points(slide17, ref_points, font_size=14)

# Save presentation
output_path = "Zoya_Ibrahim_Major_Project_PPT.pptx"
prs.save(output_path)
print(f"Presentation saved successfully at {output_path}")
